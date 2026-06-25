"""
J.A.R.V.I.S — Office Document Writer

Real .docx / .xlsx / .pptx generation. The user reported JARVIS kept
creating .md files when they asked for "Word file" — the cause was that
the only file-write tool was generic plain-text. This module gives
JARVIS proper Office output paths.

Each writer:
  - Accepts content as either plain text OR markdown
  - Converts markdown headings, bullets, bold, italic, code blocks into
    real Word formatting (not just text with # symbols)
  - Saves to a user-chosen path or auto-picks Desktop
  - Returns {success, path, size_bytes} so caller can confirm

Usage from executor:
    from core.office_writer import write_docx, write_xlsx, write_pptx
    result = write_docx(content="# Title\n\nHello world.", filename="notes.docx")
"""
from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger("jarvis.office_writer")


def _resolve_path(filename: str, ext: str) -> Path:
    """Pick a sensible default path if user gave just a filename.

    Order: explicit path → user's Desktop → user's Documents → home.
    Adds the extension if missing.
    """
    if not filename:
        filename = "jarvis_document"

    # If user gave a full path, honor it (after extension fix)
    p = Path(filename).expanduser()
    if p.is_absolute() or len(p.parts) > 1:
        if p.suffix.lower() != ext.lower():
            p = p.with_suffix(ext)
        p.parent.mkdir(parents=True, exist_ok=True)
        return p

    # Just a filename — pick default location
    if p.suffix.lower() != ext.lower():
        p = p.with_suffix(ext)
    for candidate_dir in (Path.home() / "Desktop", Path.home() / "Documents", Path.home()):
        if candidate_dir.exists():
            return candidate_dir / p.name
    return Path.cwd() / p.name


# ═══════════════════════════════════════════════════════════════════════
#  WORD (.docx)
# ═══════════════════════════════════════════════════════════════════════

def write_docx(content: str, filename: str = "jarvis_document.docx") -> dict:
    """Write content to a real Word .docx file.

    Markdown-aware: `# Heading` becomes a Word Heading, `**bold**`
    becomes bold formatting, `- item` becomes a bullet list, fenced
    code blocks become monospace paragraphs.

    Args:
        content: Markdown or plain text to write
        filename: Output filename (extension auto-added if missing)

    Returns:
        {success, path, size_bytes, paragraphs, error}
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        return {"success": False, "error": "python-docx not installed. Run: pip install python-docx"}

    out_path = _resolve_path(filename, ".docx")
    doc = Document()

    # Base style: readable body text
    try:
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
    except Exception:
        pass

    paragraph_count = 0
    in_code_block = False
    code_lines: list[str] = []
    in_bullet_list = False

    for raw_line in (content or "").split("\n"):
        line = raw_line.rstrip()

        # ── Fenced code block ────────────────────────────────────────
        if line.startswith("```"):
            if in_code_block:
                # close the code block
                code_para = doc.add_paragraph()
                code_run = code_para.add_run("\n".join(code_lines))
                code_run.font.name = "Consolas"
                code_run.font.size = Pt(9)
                try:
                    code_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                except Exception:
                    pass
                code_lines = []
                in_code_block = False
                paragraph_count += 1
            else:
                in_code_block = True
            continue
        if in_code_block:
            code_lines.append(raw_line)
            continue

        # ── Headings ────────────────────────────────────────────────
        m = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if m:
            level = min(len(m.group(1)), 4)
            doc.add_heading(m.group(2), level=level)
            paragraph_count += 1
            in_bullet_list = False
            continue

        # ── Bullet list ─────────────────────────────────────────────
        bm = re.match(r"^\s*[-*+]\s+(.+?)\s*$", line)
        if bm:
            p = doc.add_paragraph(style="List Bullet")
            _add_runs_with_formatting(p, bm.group(1))
            paragraph_count += 1
            in_bullet_list = True
            continue

        # ── Numbered list ───────────────────────────────────────────
        nm = re.match(r"^\s*\d+\.\s+(.+?)\s*$", line)
        if nm:
            p = doc.add_paragraph(style="List Number")
            _add_runs_with_formatting(p, nm.group(1))
            paragraph_count += 1
            in_bullet_list = True
            continue

        # ── Horizontal rule ─────────────────────────────────────────
        if re.match(r"^---+\s*$", line) or re.match(r"^===+\s*$", line):
            doc.add_paragraph("_" * 40)
            paragraph_count += 1
            in_bullet_list = False
            continue

        # ── Blank line ──────────────────────────────────────────────
        if not line:
            in_bullet_list = False
            continue

        # ── Regular paragraph with inline formatting ────────────────
        p = doc.add_paragraph()
        _add_runs_with_formatting(p, line)
        paragraph_count += 1
        in_bullet_list = False

    # Flush any remaining code block
    if in_code_block and code_lines:
        code_para = doc.add_paragraph()
        code_run = code_para.add_run("\n".join(code_lines))
        code_run.font.name = "Consolas"
        code_run.font.size = Pt(9)

    try:
        doc.save(str(out_path))
        size = out_path.stat().st_size
        logger.info("Wrote docx: %s (%d bytes, %d paragraphs)", out_path, size, paragraph_count)
        return {
            "success": True,
            "path": str(out_path),
            "size_bytes": size,
            "paragraphs": paragraph_count,
        }
    except Exception as e:
        return {"success": False, "error": f"Failed to save docx: {e}", "path": str(out_path)}


def _add_runs_with_formatting(paragraph, text: str):
    """Parse inline **bold** / *italic* / `code` and add runs accordingly."""
    # Split keeping the formatting markers
    parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith("*") and part.endswith("*"):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            try:
                from docx.shared import Pt
                run.font.name = "Consolas"
                run.font.size = Pt(10)
            except Exception:
                pass
        else:
            paragraph.add_run(part)


# ═══════════════════════════════════════════════════════════════════════
#  EXCEL (.xlsx)
# ═══════════════════════════════════════════════════════════════════════

def write_xlsx(content: str, filename: str = "jarvis_sheet.xlsx",
               sheet_name: str = "Sheet1") -> dict:
    """Write content to a real Excel .xlsx file.

    Accepts content as:
      - CSV text (comma-separated)
      - TSV text (tab-separated)
      - Markdown table (auto-parsed)
      - Plain text (one cell per line, single column)

    Returns: {success, path, size_bytes, rows, cols, error}
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment
    except ImportError:
        return {"success": False, "error": "openpyxl not installed. Run: pip install openpyxl"}

    out_path = _resolve_path(filename, ".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name[:31] or "Sheet1"

    rows_parsed = _parse_tabular(content or "")
    if not rows_parsed:
        rows_parsed = [["(empty)"]]

    # Write rows + auto-bold header
    for r_idx, row in enumerate(rows_parsed, 1):
        for c_idx, val in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx, value=_coerce_value(val))
            if r_idx == 1:
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center")

    # Auto-width columns (capped)
    for col_idx in range(1, max(len(r) for r in rows_parsed) + 1):
        col_letter = ws.cell(row=1, column=col_idx).column_letter
        max_len = max(
            (len(str(r[col_idx - 1])) for r in rows_parsed if col_idx - 1 < len(r)),
            default=10,
        )
        ws.column_dimensions[col_letter].width = min(max(max_len + 2, 10), 50)

    try:
        wb.save(str(out_path))
        size = out_path.stat().st_size
        logger.info("Wrote xlsx: %s (%d bytes, %d rows)", out_path, size, len(rows_parsed))
        return {
            "success": True,
            "path": str(out_path),
            "size_bytes": size,
            "rows": len(rows_parsed),
            "cols": max(len(r) for r in rows_parsed) if rows_parsed else 0,
        }
    except Exception as e:
        return {"success": False, "error": f"Failed to save xlsx: {e}", "path": str(out_path)}


def _parse_tabular(text: str) -> list[list[str]]:
    """Detect CSV / TSV / markdown table / plain text and return rows."""
    text = (text or "").strip()
    if not text:
        return []

    lines = [ln for ln in text.split("\n") if ln.strip()]

    # Markdown table: rows start/end with | and have | separators
    if lines and lines[0].strip().startswith("|") and "|" in lines[0]:
        rows = []
        for ln in lines:
            ln = ln.strip()
            if not ln.startswith("|"):
                continue
            # Skip separator row (--- | ---)
            cells = [c.strip() for c in ln.strip("|").split("|")]
            if all(re.match(r"^:?-+:?$", c) for c in cells):
                continue
            rows.append(cells)
        return rows

    # TSV
    if "\t" in lines[0]:
        return [ln.split("\t") for ln in lines]

    # CSV (look for commas in a structured way)
    if "," in lines[0] and lines[0].count(",") >= 1:
        import csv, io as _io
        reader = csv.reader(_io.StringIO(text))
        return [row for row in reader if row]

    # Fallback: one cell per line
    return [[ln] for ln in lines]


def _coerce_value(v: str):
    """Convert string to number if it looks like one — better Excel UX."""
    if v is None:
        return ""
    s = str(v).strip()
    if not s:
        return ""
    # Number?
    try:
        if "." in s or "e" in s.lower():
            return float(s)
        return int(s)
    except ValueError:
        return s


# ═══════════════════════════════════════════════════════════════════════
#  POWERPOINT (.pptx)
# ═══════════════════════════════════════════════════════════════════════

def write_pptx(content: str, filename: str = "jarvis_slides.pptx") -> dict:
    """Write content to a real PowerPoint .pptx file.

    Each '# Heading' in the content becomes a new slide title.
    Content between headings becomes bullet points on that slide.

    Returns: {success, path, size_bytes, slides, error}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"success": False, "error": "python-pptx not installed. Run: pip install python-pptx"}

    out_path = _resolve_path(filename, ".pptx")
    prs = Presentation()

    # Parse content into (title, body_bullets) slides
    slides: list[tuple[str, list[str]]] = []
    cur_title = ""
    cur_bullets: list[str] = []

    for raw_line in (content or "").split("\n"):
        line = raw_line.rstrip()
        m = re.match(r"^(#{1,3})\s+(.+?)\s*$", line)
        if m:
            if cur_title or cur_bullets:
                slides.append((cur_title or "Slide", cur_bullets))
            cur_title = m.group(2)
            cur_bullets = []
            continue
        bm = re.match(r"^\s*[-*+]\s+(.+?)\s*$", line)
        if bm:
            cur_bullets.append(bm.group(1))
            continue
        if line.strip():
            cur_bullets.append(line.strip())

    if cur_title or cur_bullets:
        slides.append((cur_title or "Slide", cur_bullets))

    if not slides:
        slides = [("JARVIS Slides", ["(no content provided)"])]

    # First slide: title layout
    title_layout = prs.slide_layouts[0]  # Title slide
    content_layout = prs.slide_layouts[1]  # Title + Content

    first_slide = prs.slides.add_slide(title_layout)
    first_slide.shapes.title.text = slides[0][0]
    if len(first_slide.placeholders) > 1 and slides[0][1]:
        first_slide.placeholders[1].text = "\n".join(slides[0][1])

    # Remaining slides: title + bullet content
    for title, bullets in slides[1:]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0] if bullets else ""
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b

    try:
        prs.save(str(out_path))
        size = out_path.stat().st_size
        logger.info("Wrote pptx: %s (%d bytes, %d slides)", out_path, size, len(slides))
        return {
            "success": True,
            "path": str(out_path),
            "size_bytes": size,
            "slides": len(slides),
        }
    except Exception as e:
        return {"success": False, "error": f"Failed to save pptx: {e}", "path": str(out_path)}


# ═══════════════════════════════════════════════════════════════════════
#  DISPATCHER — pick the right writer from filename
# ═══════════════════════════════════════════════════════════════════════

def write_document(content: str, filename: str = "") -> dict:
    """Auto-pick the right writer based on filename extension.

    Examples:
        write_document("...", "notes.docx") → write_docx
        write_document("...", "data.xlsx")  → write_xlsx
        write_document("...", "deck.pptx")  → write_pptx
    """
    fname = filename.lower()
    if fname.endswith(".docx") or fname.endswith(".doc") or "word" in fname:
        return write_docx(content, filename or "jarvis_document.docx")
    if fname.endswith(".xlsx") or fname.endswith(".xls") or "excel" in fname:
        return write_xlsx(content, filename or "jarvis_sheet.xlsx")
    if fname.endswith(".pptx") or fname.endswith(".ppt") or "power" in fname:
        return write_pptx(content, filename or "jarvis_slides.pptx")
    # Default: docx (most common request)
    return write_docx(content, filename or "jarvis_document.docx")
